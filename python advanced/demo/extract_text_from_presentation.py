from pptx import Presentation

def extract_text_from_pptx(pptx_file_path):
    presentation = Presentation(pptx_file_path)
    text = ""

    for slide_number, slide in enumerate(presentation.slides):
        text += f"Slide {slide_number + 1}:\n"

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text

if __name__ == "__main__":
    pptx_file_path = ""  # Replace with your file path
    extracted_text = extract_text_from_pptx(pptx_file_path)

    with open("output.txt", "w", encoding="utf-8") as output_file:
        output_file.write(extracted_text)

    print("Text extracted and saved to 'output.txt'")
